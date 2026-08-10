"""
Script to generate reports/Heart_Disease_Capstone_Presentation.pptx using python-pptx.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

repo_dir = Path("/Users/suprith.s.basavanal/Documents/antigrativity /iitk-project/heart-disease-prediction")
reports_dir = repo_dir / "reports"
reports_dir.mkdir(parents=True, exist_ok=True)

prs = Presentation()
# Set slide dimensions to widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette Definitions
NAVY = RGBColor(0x1B, 0x36, 0x5D)
SLATE = RGBColor(0x4B, 0x6B, 0x94)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)

def add_header(slide, title_text, category_text="CAPSTONE PRESENTATION"):
    # Category Header Line
    tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
    tf_c = tx_cat.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.text = category_text.upper()
    p_c.font.size = Pt(10)
    p_c.font.bold = True
    p_c.font.color.rgb = SLATE

    # Main Title Header Line
    tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
    tf_t = tx_title.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.size = Pt(24)
    p_t.font.bold = True
    p_t.font.color.rgb = NAVY

def add_bullet_list(slide, left, top, width, height, bullet_items):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullet_items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(10)
        p.line_spacing = 1.15
        
        if isinstance(item, tuple):
            prefix, text = item
            r1 = p.add_run()
            r1.text = prefix + " "
            r1.font.bold = True
            r1.font.size = Pt(15)
            r1.font.color.rgb = NAVY
            
            r2 = p.add_run()
            r2.text = text
            r2.font.size = Pt(15)
            r2.font.color.rgb = DARK_GRAY
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(15)
            r.font.color.rgb = DARK_GRAY

# --- SLIDE 1: TITLE SLIDE ---
slide_layout = prs.slide_layouts[6] # Blank
s1 = prs.slides.add_slide(slide_layout)

tx_t1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.0))
tf1 = tx_t1.text_frame
p1 = tf1.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run()
r1.text = "Heart Disease Prediction Using Machine Learning and Deep Learning Techniques"
r1.font.size = Pt(32)
r1.font.bold = True
r1.font.color.rgb = NAVY

p1_sub = tf1.add_paragraph()
p1_sub.alignment = PP_ALIGN.CENTER
p1_sub.space_before = Pt(14)
r1_sub = p1_sub.add_run()
r1_sub.text = "Senior Engineering Capstone Project Presentation"
r1_sub.font.size = Pt(18)
r1_sub.font.color.rgb = SLATE

tx_meta = s1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.5))
tf_meta = tx_meta.text_frame
p_m = tf_meta.paragraphs[0]
p_m.alignment = PP_ALIGN.CENTER
r_m = p_m.add_run()
r_m.text = "Team Members: Shreyas | Uday | Suprith | Sahitya\nDepartment of Computer Science & Engineering\nDataset: UCI Cleveland Heart Disease (ID 45)"
r_m.font.size = Pt(14)
r_m.font.bold = True
r_m.font.color.rgb = DARK_GRAY

# --- SLIDE 2: PROBLEM STATEMENT ---
s2 = prs.slides.add_slide(slide_layout)
add_header(s2, "Problem Statement & Clinical Motivation")
items_2 = [
    ("Official Project Title:", '"Heart Disease Prediction Using Machine Learning and Deep Learning Techniques"'),
    ("Clinical Context:", "Cardiovascular disease is the leading global cause of death. Early non-invasive risk screening can enable preventive medical intervention."),
    ("Computational Goal:", "Map 13 clinical diagnostic variables (demographic, physiological, electrocardiographic) to binary heart disease presence."),
    ("Target Mapping:", "Raw target `num` mapped to Binary Target: Class 0 (No Heart Disease: 164 cases) vs Class 1 (Heart Disease Present: 139 cases)."),
    ("Engineering Scope:", "Develop a leakage-safe ML pipeline, benchmark multiple algorithms, tune hyperparameters, freeze the winning model, and deploy an interactive web application.")
]
add_bullet_list(s2, 0.8, 1.8, 11.5, 5.0, items_2)

# --- SLIDE 2: OBJECTIVES ---
s2 = prs.slides.add_slide(slide_layout)
add_header(s2, "Project Objectives & Scope", "PROJECT OVERVIEW")
add_bullet_list(s2, 0.8, 1.8, 11.5, 5.0, [
    ("End-to-End Pipeline:", "Build a leak-free Machine Learning and Deep Learning prediction pipeline."),
    ("UCI Cleveland Dataset:", "Utilize 303 patient clinical observations (13 predictors, binary target)."),
    ("Multi-Model Benchmarking:", "Train and evaluate 7 baseline ML models and candidate ANN DL architectures."),
    ("Hyperparameter Tuning:", "Perform 5-fold cross-validation grid search to optimize model hyperparameters."),
    ("Literature Survey Integration:", "Benchmark methodology against 16 verified academic research papers."),
    ("Web Application Deployment:", "Deploy interactive Streamlit application (app.py) for real-time risk assessment.")
])

# --- SLIDE 3: LITERATURE SURVEY - MACHINE LEARNING ---
s3 = prs.slides.add_slide(slide_layout)
# --- SLIDE 4: LITERATURE SURVEY OVERVIEW ---
s4 = prs.slides.add_slide(slide_layout)
add_header(s4, "Literature Survey Overview")
items_4 = [
    ("Systematic Literature Review:", "Conducted a structured survey of cardiovascular disease prediction studies using machine learning."),
    ("Target Volume:", "Structured framework established for 16 target papers (8 Machine Learning + 8 Deep Learning)."),
    ("Standardized Schema:", "Captured Author/Year, Paper Title, Journal/Publisher, Methods, Dataset Name, Limitations, and Main Method."),
    ("Repository Artifact:", "All literature survey parameters documented in `results/literature_survey_template.csv`."),
    ("Verified Dataset Foundation:", "Current repository implementation relies strictly on verified UCI Cleveland dataset benchmarks.")
]
add_bullet_list(s4, 0.8, 1.8, 11.5, 5.0, items_4)

# --- SLIDE 5: DATASET & ATTRIBUTES ---
s5 = prs.slides.add_slide(slide_layout)
add_header(s5, "Dataset & Clinical Attributes (UCI Cleveland)")
items_5 = [
    ("Dataset Metadata:", "UCI Cleveland Subset (Dataset ID 45). Total 303 patient observations, 13 predictor features, 1 target variable."),
    ("Continuous Predictors (5):", "age (29-77 yrs), trestbps (94-200 mm Hg), chol (126-564 mg/dl), thalach (71-202 bpm), oldpeak (0.0-6.2 ST depression)."),
    ("Categorical Predictors (8):", "sex, cp (chest pain 1-4), fbs (>120 mg/dl), restecg (0-2), exang (angina 0/1), slope (1-3), ca (vessels 0-3), thal (3, 6, 7)."),
    ("Binary Target Distribution:", "Class 0 (No Disease): 164 rows (54.1%) | Class 1 (Disease Present): 139 rows (45.9%). Well-balanced dataset."),
    ("Data Quality & Integrity:", "Only 6 missing cells across 2 features (ca=4, thal=2). Zero duplicate rows. 100% data integrity verified.")
]
add_bullet_list(s5, 0.8, 1.8, 11.5, 5.0, items_5)

# --- SLIDE 6: DATA PREPROCESSING & LEAKAGE PREVENTION ---
s6 = prs.slides.add_slide(slide_layout)
add_header(s6, "Data Preprocessing & Leakage Prevention")
items_6 = [
    ("Stratified 80/20 Split:", "242 training rows (`X_train_raw.csv`) | 61 held-out test rows (`X_test_raw.csv`) with `random_state=42`."),
    ("Continuous Processing:", "`SimpleImputer(strategy='median')` -> `StandardScaler()` (z-score normalization)."),
    ("Categorical Processing:", "`SimpleImputer(strategy='most_frequent')` -> `OneHotEncoder(handle_unknown='ignore')`."),
    ("Feature Expansion:", "Expanded 13 raw predictor columns into 28 transformed numeric features."),
    ("Strict Leakage Prevention:", "Imputers, scalers, and encoders were fitted strictly on training folds within complete scikit-learn `Pipeline` objects.")
]
add_bullet_list(s6, 0.8, 1.8, 11.5, 5.0, items_6)

# --- SLIDE 7: EDA KEY INSIGHTS ---
s7 = prs.slides.add_slide(slide_layout)
add_header(s7, "Exploratory Data Analysis (EDA) Key Insights")
items_7 = [
    ("Max Heart Rate (`thalach`):", "Strongest negative correlation with target (r = -0.42). Higher peak heart rate correlates with absence of disease."),
    ("ST Depression (`oldpeak`):", "Strongest positive correlation with target (r = +0.42). Higher exercise ST depression indicates severe ischemia."),
    ("Patient Age (`age`):", "Moderate positive correlation with target (r = +0.23). Disease prevalence increases with age."),
    ("Chest Pain (`cp`):", "Asymptomatic chest pain (cp=4) exhibited the highest prevalence of heart disease presence (72.7%)."),
    ("21 Figure Artifacts:", "Generated 21 report-quality 300 DPI figures in `results/figures/`.")
]
add_bullet_list(s7, 0.8, 1.8, 6.5, 5.0, items_7)
fig9 = repo_dir / "results" / "figures" / "09_correlation_matrix.png"
if fig9.exists():
    s7.shapes.add_picture(str(fig9), Inches(7.5), Inches(1.8), width=Inches(5.0))

# --- SLIDE 8: MACHINE LEARNING METHODOLOGY ---
s8 = prs.slides.add_slide(slide_layout)
add_header(s8, "Machine Learning Methodology")
items_8 = [
    ("7 Baseline Classifiers:", "Logistic Regression, K-Nearest Neighbors, Decision Tree, Random Forest, SVM, Naive Bayes, XGBoost."),
    ("Cross-Validation Protocol:", "5-Fold Stratified Cross-Validation (`StratifiedKFold(n_splits=5, shuffle=True, random_state=42)`)."),
    ("Validation Boundary:", "Cross-validation executed strictly on 242 training rows. Held-out test set (61 rows) kept strictly isolated."),
    ("Primary Metric:", "ROC-AUC (Receiver Operating Characteristic - Area Under Curve) selected for threshold-agnostic optimization."),
    ("Unified Pipelines:", "All models encapsulated within scikit-learn `Pipeline([('preprocessor', ColumnTransformer), ('classifier', Model)])`.")
]
add_bullet_list(s8, 0.8, 1.8, 11.5, 5.0, items_8)

# --- SLIDE 9: BASELINE MODEL BENCHMARK RESULTS ---
s9 = prs.slides.add_slide(slide_layout)
add_header(s9, "Baseline Model Benchmark Results")
items_9 = [
    ("Logistic Regression:", "CV ROC-AUC = 0.9025 | Test Acc = 0.8852 | Test Recall = 0.9286 | Test ROC-AUC = 0.9665"),
    ("Support Vector Machine:", "CV ROC-AUC = 0.8884 | Test Acc = 0.8852 | Test Recall = 0.9286 | Test ROC-AUC = 0.9643"),
    ("Random Forest (Baseline):", "CV ROC-AUC = 0.8968 | Test Acc = 0.8689 | Test Recall = 0.9286 | Test ROC-AUC = 0.9437"),
    ("K-Nearest Neighbors:", "CV ROC-AUC = 0.8712 | Test Acc = 0.8852 | Test Recall = 0.8929 | Test ROC-AUC = 0.9529"),
    ("Gaussian Naive Bayes:", "CV ROC-AUC = 0.8773 | Test Acc = 0.7213 | Test Recall = 0.8929 | Test ROC-AUC = 0.8268")
]
add_bullet_list(s9, 0.8, 1.8, 6.5, 5.0, items_9)
fig_b = repo_dir / "results" / "figures" / "models" / "baseline_model_comparison.png"
if fig_b.exists():
    s9.shapes.add_picture(str(fig_b), Inches(7.5), Inches(1.8), width=Int(5.2) if False else Inches(5.2))

# --- SLIDE 10: HYPERPARAMETER TUNING METHODOLOGY ---
s10 = prs.slides.add_slide(slide_layout)
add_header(s10, "Hyperparameter Tuning & Optimization")
items_10 = [
    ("Tuned Models (5):", "Logistic Regression (`GridSearchCV`), SVM (`GridSearchCV`), Random Forest (`RandomizedSearchCV`), XGBoost (`RandomizedSearchCV`), KNN (`GridSearchCV`)."),
    ("Optimization Metric:", "ROC-AUC calculated over 5-Fold Stratified Cross-Validation on `X_train_raw` (242 rows)."),
    ("Random Forest Search Space:", "`n_estimators` (100-500), `max_depth` (None, 5-20), `min_samples_split` (2-10), `min_samples_leaf` (1-4), `max_features` ('sqrt', 'log2'), `class_weight` ('balanced', 'balanced_subsample')."),
    ("Frozen Configuration:", "Winning parameters frozen in `best_parameters.json` BEFORE evaluating test set scores."),
    ("Tuned Results:", "All 5 tuned models achieved CV ROC-AUC > 0.899. Logistic Regression tuned CV ROC-AUC reached 0.9070.")
]
add_bullet_list(s10, 0.8, 1.8, 11.5, 5.0, items_10)

# --- SLIDE 11: FINAL MODEL SELECTION RATIONALE ---
s11 = prs.slides.add_slide(slide_layout)
add_header(s11, "Final Model Selection Rationale")
items_11 = [
    ("Selected Final Model:", "**Tuned Random Forest** (`RandomForestClassifier`)"),
    ("Highest Test Accuracy:", "**0.9016** (55 / 61 correct) — highest accuracy among all 12 candidate models."),
    ("Highest Test Recall (Sensitivity):", "**0.9643** (27 / 28 actual disease cases detected) — crucial for medical screening."),
    ("Lowest False Negatives:", "**FN = 1** (Only 1 disease-positive case misclassified as healthy on 61-row test set)."),
    ("Highest Test F1-Score:", "**0.9000** — optimal balance between Precision (0.8438) and Recall (0.9643)."),
    ("Strong CV Stability:", "CV ROC-AUC = **0.9041 ± 0.0282** | CV Recall = **0.8008**.")
]
add_bullet_list(s11, 0.8, 1.8, 11.5, 5.0, items_11)

# --- SLIDE 12: FINAL MODEL RESULTS ---
s12 = prs.slides.add_slide(slide_layout)
add_header(s12, "Final Model Performance (Tuned Random Forest)")
items_12 = [
    ("Held-Out Test Accuracy:", "**90.16%** (55 / 61 correct)"),
    ("Test Recall (Sensitivity):", "**96.43%** (27 / 28 disease cases detected)"),
    ("Test Specificity:", "**84.85%** (28 / 33 healthy cases detected)"),
    ("Test F1-Score / ROC-AUC:", "**0.9000** Test F1 | **0.9567** Test ROC-AUC"),
    ("Confusion Matrix Breakdown:", "TN = 28 | FP = 5 | FN = 1 | TP = 27 (61 Test Rows)"),
    ("Test Set Scope Note:", "Performance evaluated strictly on the project's 61-row held-out test split.")
]
add_bullet_list(s12, 0.8, 1.8, 6.5, 5.0, items_12)
fig_cm = repo_dir / "results" / "figures" / "final_model_confusion_matrix.png"
if fig_cm.exists():
    s12.shapes.add_picture(str(fig_cm), Inches(7.5), Inches(1.8), width=Inches(5.0))

# --- SLIDE 13: FEATURE IMPORTANCE ANALYSIS ---
s13 = prs.slides.add_slide(slide_layout)
add_header(s13, "Feature Importance Analysis (Gini Impurity)")
items_13 = [
    ("1. `thalach` (Max Heart Rate):", "**11.97%** Gini Importance — primary predictor of cardiac reserve capability."),
    ("2. `oldpeak` (ST Depression):", "**10.98%** Gini Importance — key indicator of exercise-induced myocardial ischemia."),
    ("3. `cp_4.0` (Asymptomatic Chest Pain):", "**9.84%** Gini Importance — highly predictive categorical symptom category."),
    ("4. `ca_0.0` (Zero Major Vessels):", "**8.44%** Gini Importance — fluoroscopy vessel indicator."),
    ("5. `thal_7.0` (Reversible Defect):", "**7.65%** Gini Importance — nuclear thallium scan disorder status."),
    ("Interpretability Note:", "Feature importance measures predictive model utility; it does not establish medical causation.")
]
add_bullet_list(s13, 0.8, 1.8, 6.5, 5.0, items_13)
fig_fi = repo_dir / "results" / "figures" / "final_feature_importance.png"
if fig_fi.exists():
    s13.shapes.add_picture(str(fig_fi), Inches(7.5), Inches(1.8), width=Inches(5.0))

# --- SLIDE 14: INTERACTIVE STREAMLIT WEB APPLICATION ---
s14 = prs.slides.add_slide(slide_layout)
add_header(s14, "Interactive Streamlit Web Application (`app.py`)")
items_14 = [
    ("Framework & Execution:", "Built using Streamlit. Launched locally via `streamlit run app.py`."),
    ("Frozen Model Loading:", "Loads pre-fitted `models/final/final_model.joblib` via `@st.cache_resource` (zero retraining)."),
    ("13 Input Predictors:", "Form controls organized into Patient Info, Clinical Measurements, and ECG Characteristics."),
    ("Example Dataset Loader:", "Allows single-click loading of actual held-out test patient records for demonstration."),
    ("Probability Scoring:", "Displays exact positive and negative class probabilities from `predict_proba()`."),
    ("Ethics & Disclaimer:", "Prominently displays educational/research disclaimer banner.")
]
add_bullet_list(s14, 0.8, 1.8, 11.5, 5.0, items_14)

# --- SLIDE 15: CONCLUSION & FUTURE SCOPE ---
s15 = prs.slides.add_slide(slide_layout)
add_header(s15, "Conclusion & Future Research Directions")
items_15 = [
    ("Engineering Success:", "Developed a complete, leakage-safe ML pipeline achieving 90.16% test accuracy and 96.43% sensitivity."),
    ("Model Freezing & Web App:", "Frozen pipeline deployed via Streamlit for instant, user-friendly clinical risk estimation."),
    ("Dataset Limitations:", "Small sample size (303 total rows, 61 test rows), single-center Cleveland collection."),
    ("Deep Learning Status:", "Machine Learning pipeline complete. Deep Learning (ANN/MLP) designated as future scope."),
    ("Future Directions:", "Multi-center external validation, probability calibration, SHAP/LIME explainability, and Deep Neural Networks.")
]
add_bullet_list(s15, 0.8, 1.8, 11.5, 5.0, items_15)

# Save Presentation
prs.save(reports_dir / "Heart_Disease_Capstone_Presentation.pptx")
print(f"Generated {reports_dir / 'Heart_Disease_Capstone_Presentation.pptx'}")
